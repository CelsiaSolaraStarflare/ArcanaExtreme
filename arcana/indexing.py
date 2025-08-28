import os
import pandas as pd
from docx import Document
from pptx import Presentation
import chardet
from arcana.fiber import FiberDBMS
import nltk

# Ensure NLTK data is available before importing NLTK functions
import arcana.nltk_setup
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
from PyPDF2 import PdfReader
import csv
import re
import ast
from scripts.config import INDEX_FILE

# NLTK data is now downloaded once in Arcanalte.py at startup.

def extract_keywords(text, lang='en'):
    # Handles both English and Chinese, and ignores emojis/symbols
    stop_words = set(stopwords.words('english')) if lang == 'en' else set()
    # Tokenize Chinese and English, keep CJK, ignore symbols/emojis
    words = re.findall(r'[\u4e00-\u9fff]+|[a-zA-Z]+', text)
    if lang == 'en':
        return [w for w in words if w.lower() not in stop_words]
    return words

def detect_language(text):
    # Simple check: if contains CJK, treat as Chinese
    if re.search(r'[\u4e00-\u9fff]', text):
        return 'zh'
    return 'en'

def indexing(cache_dir: str):
    """
    Traverses a directory, processes all supported files, extracts content and
    keywords, and builds a search index using FiberDBMS.

    Args:
        cache_dir (str): The path to the directory containing files to be indexed.

    Returns:
        int: The total number of entries indexed.
    """
    dbms = FiberDBMS()
    # Load existing database if present to avoid duplicates
    existing_entries = set()
    if os.path.exists(INDEX_FILE):
        try:
            dbms.load_from_file(INDEX_FILE)
            existing_entries = {(e['name'], e['content']) for e in dbms.database}
            print(f"Loaded existing index with {len(existing_entries)} entries. New indexing will skip duplicates.")
        except Exception as exc:
            print(f"Could not load existing index for duplicate checking: {exc}")

    entries = []

    # Traverse the cache directory for all supported file types
    for root, _, files in os.walk(cache_dir):
        for file in files:
            file_path = os.path.join(root, file)
            file_extension = os.path.splitext(file)[1].lower()

            try:
                # Process different file types and extract content
                if file_extension == ".txt":
                    with open(file_path, 'rb') as f:
                        raw_data = f.read()
                        encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'
                    with open(file_path, 'r', encoding=encoding, errors='replace') as f:
                        content = f.read()
                elif file_extension == ".docx":
                    doc = Document(file_path)
                    content = "\n".join([para.text for para in doc.paragraphs])
                elif file_extension == ".pptx":
                    presentation = Presentation(file_path)
                    all_texts = []
                    for slide in presentation.slides:
                        slide_texts = []
                        if slide.shapes.title:
                            slide_texts.append(slide.shapes.title.text)
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                slide_texts.append(shape.text_frame.text)  # type: ignore
                        all_texts.append("\n".join(slide_texts))
                    content = "\n".join(all_texts)
                elif file_extension in [".xls", ".xlsx", ".csv"]:
                    df = pd.read_excel(file_path) if "xls" in file_extension else pd.read_csv(file_path)
                    content = df.to_csv(index=False)
                elif file_extension == ".pdf":
                    reader = PdfReader(file_path)
                    content = "\n".join([page.extract_text() or '' for page in reader.pages])
                else:
                    content = None  # Ignore unsupported formats

                # Add the file content to the database
                if content:  # Ensure content is not None
                    for i in content.split('\n'):
                        i = i.strip()
                        if i and (file, i) not in existing_entries:
                            lang = detect_language(i)
                            keywords = extract_keywords(i, lang)
                            entries.append([file, i, ','.join(keywords)])
                            existing_entries.add((file, i))  # avoid duplicates within same run
                #print('The database has been indexed')
                print(f"Processed {file}: {len(entries)} entries indexed.")
            except Exception as e:
                print(f"Failed to process {file}: {e}")
    print(f"Indexed {len(entries)} entries from {cache_dir}")

    # Add all new entries to dbms in one go
    for name, content, tags in entries:
        dbms.add_entry(name=name, content=content, tags=tags.split(','))

    # Save the database using the dbms's save method to the configured file
    dbms.save(INDEX_FILE)
    print(f"Database saved to {INDEX_FILE}")
    return len(entries)

def correct_malformed_row(row):
    # If row is a dict with a single key, try to split it
    if isinstance(row, dict) and len(row) == 1:
        key = list(row.keys())[0]
        fields = key.split('\t')
        if len(fields) == 4:
            name, timestamp, content, tags = fields
            if tags.startswith('[') and tags.endswith(']'):
                try:
                    tags_list = ast.literal_eval(tags)
                    tags = ','.join(str(t).strip() for t in tags_list)
                except Exception:
                    tags = tags
            return {
                'name': name,
                'timestamp': timestamp,
                'content': content,
                'tags': tags
            }
    return None
