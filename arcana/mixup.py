# mypy: skip-file
# flake8: noqa
import streamlit as st
import os
import datetime
import re
import io
import requests

# Ensure NLTK data is available before importing NLTK functions
import arcana.nltk_setup
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords

from response import openai_api_call
from arcana.fiber import FiberDBMS
from scripts.config import GENERATED_FILES_DIR
from openai.types.chat import ChatCompletionMessageParam
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE
from docx import Document
from docx.shared import Pt as DocxPt, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH  # type: ignore

# --- Markdown Parser for Word Documents ---

def parse_mixed_formatting_to_runs(paragraph, text):
    """
    Parses text that may contain both markdown and custom font tags,
    and adds formatted runs to a Word paragraph.
    Handles: **bold**, *italic*, `code`, [links], [font] tags, etc.
    """
    # First handle custom font tags which have priority
    parts = []
    current_text = text
    
    # Split by custom font tags
    font_parts = re.split(r'(\[font[^\]]*\]|\[/font\])', current_text)
    style_stack = [{'bold': None, 'italic': None, 'underline': None, 'name': None, 'size': None, 'color': None}]
    
    for part in font_parts:
        if not part:
            continue
            
        if part.startswith('[font'):
            # Parse font attributes
            new_style = style_stack[-1].copy()
            attrs = {}
            attr_string = part[part.find(' '):-1].strip() if ' ' in part else ''
            pairs = re.findall(r'(\w+)=(".*?"|\'.*?\'|[^\s\']*)', attr_string)
            for key, value in pairs:
                attrs[key.lower()] = value.strip("'\"")

            if 'bold' in attrs: new_style['bold'] = attrs['bold'].lower() == 'true'
            if 'italic' in attrs: new_style['italic'] = attrs['italic'].lower() == 'true'
            if 'underline' in attrs: new_style['underline'] = attrs['underline'].lower() == 'true'
            if 'name' in attrs: new_style['name'] = attrs['name']
            if 'size' in attrs:
                try: new_style['size'] = DocxPt(int(attrs['size']))
                except ValueError: pass
            if 'color' in attrs:
                try: new_style['color'] = DocxRGBColor.from_string(attrs['color'].replace("#", ""))
                except (ValueError, KeyError): pass
            style_stack.append(new_style)
            
        elif part == '[/font]':
            if len(style_stack) > 1:
                style_stack.pop()
                
        else:
            # This is actual text content - apply current font style and then parse markdown
            current_style = style_stack[-1]
            
            # Now parse markdown within this text part
            markdown_parts = parse_markdown_inline(part)
            
            for md_part in markdown_parts:
                combined_style = current_style.copy()
                
                # Combine custom font style with markdown style
                if md_part.get('bold'): combined_style['bold'] = True
                if md_part.get('italic'): combined_style['italic'] = True
                if md_part.get('code'): 
                    combined_style['name'] = 'Consolas'
                    combined_style['background'] = 'light_gray'
                if md_part.get('underline'): combined_style['underline'] = True
                if md_part.get('link'): 
                    combined_style['color'] = DocxRGBColor(0, 0, 255)
                    combined_style['underline'] = True
                
                parts.append({
                    'text': md_part['text'], 
                    'style': combined_style
                })
    
    # Apply all parts to the paragraph
    for part in parts:
        if not part['text'].strip():
            continue
            
        run = paragraph.add_run(part['text'])
        style = part['style']
        
        # Apply formatting
        if style.get('bold'):
            run.font.bold = True
        if style.get('italic'):
            run.font.italic = True
        if style.get('underline'):
            run.font.underline = True
        if style.get('name'):
            run.font.name = style['name']
        if style.get('size'):
            run.font.size = style['size']
        if style.get('color'):
            run.font.color.rgb = style['color']
        if style.get('background') == 'light_gray':
            run.font.highlight_color = 7

def parse_markdown_inline(text):
    """
    Parses inline markdown formatting and returns a list of text parts with formatting info.
    """
    # Patterns for inline markdown
    patterns = [
        (r'\*\*(.*?)\*\*', 'bold'),      # **bold**
        (r'(?<!\*)\*([^*]+?)\*(?!\*)', 'italic'),  # *italic* (not part of **)
        (r'`(.*?)`', 'code'),            # `code`
        (r'_\_(.*?)__', 'underline'),    # __underline__
        (r'(?<!_)_([^_]+?)_(?!_)', 'italic'),      # _italic_ (not part of __)
        (r'\[(.*?)\]\((.*?)\)', 'link'), # [text](url)
    ]
    
    # Start with the full text
    parts = [{'text': text, 'formatting': {}}]
    
    # Apply each pattern
    for pattern, format_type in patterns:
        new_parts = []
        for part in parts:
            if part.get('processed'):
                new_parts.append(part)
                continue
                
            text_to_process = part['text']
            current_formatting = part.get('formatting', {}).copy()
            
            # Find all matches
            matches = list(re.finditer(pattern, text_to_process))
            if not matches:
                new_parts.append(part)
                continue
            
            last_end = 0
            for match in matches:
                # Add text before match
                if match.start() > last_end:
                    before_text = text_to_process[last_end:match.start()]
                    if before_text:
                        new_parts.append({'text': before_text, 'formatting': current_formatting.copy()})
                
                # Add formatted text
                if format_type == 'link':
                    link_text = match.group(1)
                    link_url = match.group(2)
                    link_formatting = current_formatting.copy()
                    link_formatting['link'] = link_url
                    new_parts.append({'text': link_text, 'formatting': link_formatting, 'processed': True})
                else:
                    matched_text = match.group(1)
                    matched_formatting = current_formatting.copy()
                    matched_formatting[format_type] = True
                    new_parts.append({'text': matched_text, 'formatting': matched_formatting, 'processed': True})
                
                last_end = match.end()
            
            # Add remaining text
            if last_end < len(text_to_process):
                remaining_text = text_to_process[last_end:]
                if remaining_text:
                    new_parts.append({'text': remaining_text, 'formatting': current_formatting.copy()})
        
        parts = new_parts
    
    # Convert to final format
    result = []
    for part in parts:
        formatting = part.get('formatting', {})
        result.append({
            'text': part['text'],
            'bold': formatting.get('bold', False),
            'italic': formatting.get('italic', False),
            'code': formatting.get('code', False),
            'underline': formatting.get('underline', False),
            'link': formatting.get('link', None)
        })
    
    return result

def parse_markdown_to_word_runs(paragraph, text):
    """
    Parses markdown-like text and adds formatted runs to a Word paragraph.
    Supports: **bold**, *italic*, `code`, [links], and basic formatting.
    """
    # Use the new mixed formatting parser for backward compatibility
    parse_mixed_formatting_to_runs(paragraph, text)

def parse_markdown_content_to_word(doc, content):
    """
    Parses markdown content and adds it to a Word document with proper formatting.
    Handles headers, lists, paragraphs, code blocks, tables, and inline formatting.
    """
    lines = content.split('\n')
    current_list_level = 0
    in_code_block = False
    code_block_content = []
    in_table = False
    table_rows = []
    
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped_line = line.strip()
        
        # Handle code blocks
        if stripped_line.startswith('```'):
            if in_code_block:
                # End code block
                if code_block_content:
                    code_para = doc.add_paragraph()
                    code_run = code_para.add_run('\n'.join(code_block_content))
                    code_run.font.name = 'Consolas'
                    code_run.font.size = DocxPt(10)
                    code_para.paragraph_format.left_indent = DocxPt(20)
                code_block_content = []
                in_code_block = False
            else:
                in_code_block = True
            i += 1
            continue
        
        if in_code_block:
            code_block_content.append(line)
            i += 1
            continue
        
        # Handle tables
        if '|' in stripped_line and stripped_line.count('|') >= 2:
            # Check if this looks like a table
            if not in_table:
                # Start collecting table rows
                table_rows = []
                in_table = True
            
            # Parse table row
            cells = [cell.strip() for cell in stripped_line.split('|')]
            # Remove empty cells at start/end (markdown table format)
            if cells and not cells[0]:
                cells = cells[1:]
            if cells and not cells[-1]:
                cells = cells[:-1]
            
            # Check if this is a separator row (contains only dashes, spaces, and |)
            if all(re.match(r'^[\s\-\:]*$', cell) for cell in cells):
                # This is a separator row, skip it but keep building table
                i += 1
                continue
            
            table_rows.append(cells)
            i += 1
            continue
        else:
            # End of table if we were in one
            if in_table and table_rows:
                # Create the table
                if table_rows:
                    num_cols = max(len(row) for row in table_rows) if table_rows else 1
                    table = doc.add_table(rows=len(table_rows), cols=num_cols)
                    table.style = 'Table Grid'
                    
                    for row_idx, row_data in enumerate(table_rows):
                        for col_idx, cell_text in enumerate(row_data):
                            if col_idx < num_cols:
                                cell = table.cell(row_idx, col_idx)
                                # Handle custom font tags and markdown in cells
                                parse_mixed_formatting_to_runs(cell.paragraphs[0], cell_text)
                                # Make header row bold
                                if row_idx == 0:
                                    for run in cell.paragraphs[0].runs:
                                        run.font.bold = True
                
                table_rows = []
                in_table = False
        
        # Skip empty lines
        if not stripped_line:
            i += 1
            continue
        
        # Handle headers (including custom font tags)
        if stripped_line.startswith('#'):
            header_level = len(stripped_line) - len(stripped_line.lstrip('#'))
            header_text = stripped_line.lstrip('# ').strip()
            
            if header_level <= 3:
                heading = doc.add_heading(header_text, level=header_level)
                # Handle custom font tags in headers
                if '[font' in header_text:
                    heading.clear()
                    parse_mixed_formatting_to_runs(heading, header_text)
            else:
                # For h4, h5, h6, use bold paragraph
                para = doc.add_paragraph()
                parse_mixed_formatting_to_runs(para, header_text)
                # Make it bold by default
                for run in para.runs:
                    run.font.bold = True
                    run.font.size = DocxPt(12 if header_level == 4 else 11)
        
        # Handle lists
        elif stripped_line.startswith(('- ', '* ', '+ ')):
            list_text = stripped_line[2:].strip()
            para = doc.add_paragraph(style='List Bullet')
            parse_mixed_formatting_to_runs(para, list_text)
        
        elif re.match(r'^\d+\. ', stripped_line):
            list_text = re.sub(r'^\d+\. ', '', stripped_line).strip()
            para = doc.add_paragraph(style='List Number')
            parse_mixed_formatting_to_runs(para, list_text)
        
        # Handle blockquotes
        elif stripped_line.startswith('>'):
            quote_text = stripped_line.lstrip('> ').strip()
            para = doc.add_paragraph()
            para.paragraph_format.left_indent = DocxPt(20)
            run = para.add_run(f'"{quote_text}"')
            run.font.italic = True
        
        # Handle horizontal rules
        elif stripped_line in ['---', '***', '___']:
            para = doc.add_paragraph('_' * 50)
            para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Handle regular paragraphs
        else:
            para = doc.add_paragraph()
            parse_mixed_formatting_to_runs(para, stripped_line)
        
        i += 1
    
    # Handle any remaining table at the end
    if in_table and table_rows:
        if table_rows:
            num_cols = max(len(row) for row in table_rows) if table_rows else 1
            table = doc.add_table(rows=len(table_rows), cols=num_cols)
            table.style = 'Table Grid'
            
            for row_idx, row_data in enumerate(table_rows):
                for col_idx, cell_text in enumerate(row_data):
                    if col_idx < num_cols:
                        cell = table.cell(row_idx, col_idx)
                        parse_mixed_formatting_to_runs(cell.paragraphs[0], cell_text)
                        if row_idx == 0:
                            for run in cell.paragraphs[0].runs:
                                run.font.bold = True

# --- State Management ---

def init_presentation_state(force_reset=False):
    """Initializes or resets the session state for the presentation generator."""
    if force_reset:
        st.session_state.presentation_step = "initial"
        st.session_state.presentation_topic = ""
        st.session_state.presentation_outline = ""
        st.session_state.presentation_content = []
        return

    if 'presentation_step' not in st.session_state:
        st.session_state.presentation_step = "initial"
    if 'presentation_topic' not in st.session_state:
        st.session_state.presentation_topic = ""
    if 'presentation_outline' not in st.session_state:
        st.session_state.presentation_outline = ""
    if 'presentation_content' not in st.session_state:
        st.session_state.presentation_content = []

def init_study_guide_state(force_reset=False):
    """Initializes or resets the session state for the study guide generator."""
    if force_reset:
        st.session_state.study_guide_step = "initial"
        st.session_state.study_guide_topic = ""
        st.session_state.study_guide_style = "comprehensive"
        st.session_state.study_guide_extreme_mode = False
        st.session_state.study_guide_page_count = -1
        st.session_state.study_guide_outline = ""
        st.session_state.study_guide_content = []
        return

    if 'study_guide_step' not in st.session_state:
        st.session_state.study_guide_step = "initial"
    if 'study_guide_topic' not in st.session_state:
        st.session_state.study_guide_topic = ""
    if 'study_guide_style' not in st.session_state:
        st.session_state.study_guide_style = "comprehensive"
    if 'study_guide_extreme_mode' not in st.session_state:
        st.session_state.study_guide_extreme_mode = False
    if 'study_guide_page_count' not in st.session_state:
        st.session_state.study_guide_page_count = -1
    if 'study_guide_outline' not in st.session_state:
        st.session_state.study_guide_outline = ""
    if 'study_guide_content' not in st.session_state:
        st.session_state.study_guide_content = []

# --- Helper Functions ---

def get_context_for_topic(dbms, topic):
    """Extracts keywords from a topic and queries the database for relevant context."""
    stop_words = set(stopwords.words('english'))
    words = word_tokenize(topic)
    keywords = [word for word in words if word.lower() not in stop_words and word.isalpha()]
    
    if not keywords:
        return "" # Return empty string if no keywords are found

    results = dbms.query(" ".join(keywords), top_n=10)
    
    if not results:
        return "" # Return empty string if no context is found

    context = "Here is some relevant information from your documents:\n\n"
    for result in results:
        context += f"--- Start of content from {result['name']} ---\n"
        context += f"{result['content']}\n"
        context += f"--- End of content from {result['name']} ---\n\n"
    return context

def parse_outline_to_slides(outline_text):
    """Parses a markdown-formatted outline into a list of slide dictionaries."""
    slides = []
    current_slide = None
    for line in outline_text.splitlines():
        line = line.strip()
        if line.startswith("####"):
            if current_slide:
                slides.append(current_slide)
            current_slide = {"title": line.lstrip("# ").strip(), "points": []}
        elif line.startswith("-") and current_slide is not None:
            current_slide["points"].append(line.lstrip("- ").strip())
    if current_slide:
        slides.append(current_slide)
    return slides

def parse_outline_to_sections(outline_text):
    """Parses a markdown-formatted outline into a list of study guide sections."""
    sections = []
    current_section = None
    for line in outline_text.splitlines():
        line = line.strip()
        if line.startswith("##"):  # Main sections
            if current_section:
                sections.append(current_section)
            current_section = {"title": line.lstrip("# ").strip(), "subsections": []}
        elif line.startswith("-") and current_section is not None:
            current_section["subsections"].append(line.lstrip("- ").strip())
    if current_section:
        sections.append(current_section)
    return sections

def get_study_guide_style_prompt(style):
    """Returns the appropriate prompt modification based on study guide style."""
    style_prompts = {
        "comprehensive": "Create a detailed, comprehensive study guide with thorough explanations, examples, and practice questions.",
        "summary": "Create a concise summary-style study guide focusing on key points and essential information.",
        "outline": "Create a structured outline-format study guide with hierarchical organization of topics.",
        "flashcard_prep": "Create a study guide optimized for flashcard creation with clear definitions and Q&A format.",
        "exam_focused": "Create an exam-focused study guide with emphasis on likely test questions and critical concepts."
    }
    return style_prompts.get(style, style_prompts["comprehensive"])

def get_page_count_instructions(page_count, total_sections):
    """Returns content length instructions based on target page count."""
    if page_count == -1:
        return "Write comprehensive content with natural length based on topic complexity."
    
    # Estimate words per page for study guide format (160-200 words per page)
    # Using 180 words per page as a middle ground for study guides with headers, lists, and formatting
    words_per_page = 180
    total_target_words = page_count * words_per_page
    words_per_section = total_target_words // max(total_sections, 1)
    
    if page_count <= 5:
        length_instruction = f"Keep content concise and focused. Aim for approximately {words_per_section} words per section (~{page_count} total pages at ~{words_per_page} words/page for study guide format)."
    elif page_count <= 15:
        length_instruction = f"Provide balanced detail with moderate length. Aim for approximately {words_per_section} words per section (~{page_count} total pages at ~{words_per_page} words/page for study guide format)."
    elif page_count <= 30:
        length_instruction = f"Create comprehensive, detailed content. Aim for approximately {words_per_section} words per section (~{page_count} total pages at ~{words_per_page} words/page for study guide format)."
    else:
        length_instruction = f"Generate extensive, thorough content with maximum detail. Aim for approximately {words_per_section} words per section (~{page_count} total pages at ~{words_per_page} words/page for study guide format)."
    
    # Add specific reminder about study guide formatting
    length_instruction += f" Remember to use study guide formatting with headers, bullet points, and clear sections which affects word density per page."
    
    return length_instruction

def get_markdown_formatting_instructions():
    """Returns instructions for proper markdown formatting in generated content."""
    return """
Format your response using proper markdown:
- Use **bold** for key terms and important concepts
- Use *italic* for emphasis and definitions
- Use `code formatting` for technical terms, formulas, or specific examples
- Use ### Headers for main sections and #### for subsections
- Use - bullet points for lists and key points
- Use 1. numbered lists for step-by-step processes
- Use > blockquotes for important notes or quotes
- Use --- for section dividers where appropriate
"""

def add_formatted_text_to_shape(shape, text: str):
    """
    Parses text with formatting tags and adds bullet points and alignment.
    """
    tf = shape.text_frame
    tf.clear()
    current_alignment = PP_ALIGN.LEFT
    default_size = Pt(18)
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        # Alignment tags
        if stripped.lower() == '[align=center]':
            current_alignment = PP_ALIGN.CENTER
            continue
        if stripped.lower() == '[/align]':
            current_alignment = PP_ALIGN.LEFT
            continue
        # Font size tags
        size_match = re.match(r"\[font size=(\d+)\]", stripped, re.I)
        if size_match:
            default_size = Pt(int(size_match.group(1)))
            continue
        # Clean tags
        clean_line = re.sub(r"\[[^\]]+\]", '', stripped)
        p = tf.add_paragraph()
        if clean_line.startswith('- '):
            p.text = clean_line[2:]
            p.level = 1
        else:
            p.text = clean_line
        p.alignment = current_alignment
        p.font.size = default_size
    return

# --- Document Generation ---

def create_presentation_from_content(presentation_content, topic, api_key):
    """Creates a PowerPoint presentation from a list of slide content, including images."""
    # Load template if available (located in plugin/templates/template.pptx)
    template_path = os.path.join(os.path.dirname(__file__), 'templates', 'template.pptx')
    if os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()
    # Use built-in Title and Content slide layout for consistent formatting
    content_slide_layout = prs.slide_layouts[1]

    # Create slides using the selected layout
    for idx, slide_data in enumerate(presentation_content, start=1):
        slide = prs.slides.add_slide(content_slide_layout)
        # Set slide title
        slide.shapes.title.text = slide_data['title']  # type: ignore
        title_p = slide.shapes.title.text_frame.paragraphs[0]  # type: ignore
        title_p.font.bold = True
        title_p.font.size = Pt(32)
        title_p.alignment = PP_ALIGN.CENTER

        # Parse and separate main content and speaker notes
        content_lines, notes_lines = split_content_and_notes(slide_data['content'], slide_data['title'])
        # Remove duplicate title lines
        content_lines = [l for l in content_lines if l.strip() and l.strip() != slide_data['title']]
        # Remove image tags
        content_lines = [l for l in content_lines if not l.startswith('[image]')]
        text_content = "\n".join(content_lines)
        if text_content.strip():
            # Use placeholder for content and clear existing content
            content_shape = slide.shapes.placeholders[1]  # type: ignore
            content_shape.text_frame.clear()  # type: ignore
            content_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # type: ignore
            add_formatted_text_to_shape(content_shape, text_content)
        # Add speaker notes if present
        if notes_lines:
            notes_slide = slide.notes_slide  # type: ignore
            notes_tf = notes_slide.notes_text_frame
            notes_tf.clear()
            for note_line in notes_lines:
                notes_tf.add_paragraph().text = note_line

    # Add slide numbers to each slide
    for num, slide in enumerate(prs.slides, start=1):
        number_box = slide.shapes.add_textbox(
            prs.slide_width - Inches(1.0),  # type: ignore
            prs.slide_height - Inches(0.5),  # type: ignore
            Inches(1.0),
            Inches(0.5)
        )
        num_tf = number_box.text_frame
        p = num_tf.paragraphs[0]
        p.text = str(num)
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.RIGHT

    # Save presentation
    safe_topic = "".join(c for c in topic if c.isalnum() or c in (' ', '_')).rstrip()
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"Presentation_{safe_topic}_{timestamp}.pptx"
    
    os.makedirs(GENERATED_FILES_DIR, exist_ok=True)
    file_path = os.path.join(GENERATED_FILES_DIR, filename)
    prs.save(file_path)
    return file_path

def create_document_from_content(presentation_content, topic):
    """Creates a Word document from a list of slide content."""
    doc = Document()
    doc.add_heading(f"Report on: {topic}", level=0)
    for slide_data in presentation_content:
        doc.add_heading(slide_data['title'], level=1)
        # Use the new formatter for the document content
        add_formatted_text_to_document(doc, slide_data['content'])
        doc.add_paragraph() # Add some space

    safe_topic = "".join(c for c in topic if c.isalnum() or c in (' ', '_')).rstrip()
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"Document_{safe_topic}_{timestamp}.docx"
    
    os.makedirs(GENERATED_FILES_DIR, exist_ok=True)
    file_path = os.path.join(GENERATED_FILES_DIR, filename)
    doc.save(file_path)
    return file_path

def create_study_guide_from_content(study_guide_content, topic, style, extreme_mode=False):
    """Creates a Word document study guide from a list of section content."""
    doc = Document()
    
    # Add title and metadata
    title = doc.add_heading(f"Study Guide: {topic}", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Add metadata
    metadata_p = doc.add_paragraph()
    metadata_p.add_run(f"Generated: {datetime.datetime.now().strftime('%B %d, %Y at %I:%M %p')}\n")
    metadata_p.add_run(f"Style: {style.replace('_', ' ').title()}\n")
    metadata_p.add_run(f"Generation Mode: {'üöÄ Extreme Mode' if extreme_mode else '‚ö° Fast Mode'}\n")
    metadata_p.add_run(f"Sections: {len(study_guide_content)}")
    metadata_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_page_break()
    
    # Add table of contents
    toc_heading = doc.add_heading("Table of Contents", level=1)
    for idx, section_data in enumerate(study_guide_content, 1):
        toc_p = doc.add_paragraph()
        toc_p.add_run(f"{idx}. {section_data['title']}")
        toc_p.style = 'List Number'
    
    doc.add_page_break()
    
    # Add content sections
    for section_data in study_guide_content:
        doc.add_heading(section_data['title'], level=1)
        add_formatted_text_to_document(doc, section_data['content'])
        doc.add_paragraph()  # Add spacing between sections

    # Save the document
    safe_topic = "".join(c for c in topic if c.isalnum() or c in (' ', '_')).rstrip()
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    mode_suffix = "extreme" if extreme_mode else "fast"
    filename = f"StudyGuide_{safe_topic}_{style}_{mode_suffix}_{timestamp}.docx"
    
    os.makedirs(GENERATED_FILES_DIR, exist_ok=True)
    file_path = os.path.join(GENERATED_FILES_DIR, filename)
    doc.save(file_path)
    return file_path

def add_formatted_text_to_document(doc, text):
    """
    Parses both markdown and custom-formatted text and adds it to a Word document,
    preserving formatting like bold, italics, underline, color, alignment, headers, and lists.
    Now supports both markdown syntax and custom formatting tags.
    """
    current_alignment = WD_ALIGN_PARAGRAPH.LEFT
    
    # First check if the text contains markdown patterns
    has_markdown = any([
        '**' in text,  # Bold
        re.search(r'(?<!\*)\*(?!\*)', text),  # Italic (single asterisk not part of **)
        '`' in text,   # Code
        re.search(r'^#+\s', text, re.MULTILINE),  # Headers
        re.search(r'^\s*[-*+]\s', text, re.MULTILINE),  # Lists
        re.search(r'^\s*\d+\.\s', text, re.MULTILINE),  # Numbered lists
        '___' in text or '---' in text or '***' in text,  # Horizontal rules
    ])
    
    # If contains markdown, use the new markdown parser
    if has_markdown:
        parse_markdown_content_to_word(doc, text)
        return
    
    # Otherwise, use the existing custom formatting tag parser for backward compatibility
    for line in text.splitlines():
        # Skip image tags and empty lines
        if line.startswith('[image') or not line.strip():
            continue

        align_open_match = re.match(r"^\s*\[align=(center|right|left|justify)\]\s*$", line, re.I)
        align_close_match = re.match(r"^\s*\[/align\]\s*$", line, re.I)

        if align_open_match:
            alignment_str = align_open_match.group(1).upper()
            if alignment_str == 'LEFT': current_alignment = WD_ALIGN_PARAGRAPH.LEFT
            elif alignment_str == 'CENTER': current_alignment = WD_ALIGN_PARAGRAPH.CENTER
            elif alignment_str == 'RIGHT': current_alignment = WD_ALIGN_PARAGRAPH.RIGHT
            elif alignment_str == 'JUSTIFY': current_alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            continue
        
        if align_close_match:
            current_alignment = WD_ALIGN_PARAGRAPH.LEFT
            continue

        p = doc.add_paragraph()
        p.alignment = current_alignment

        parts = re.split(r"(\[font[^\]]*\]|\[/font\])", line)
        style_stack = [{'bold': None, 'italic': None, 'underline': None, 'name': None, 'size': None, 'color': None}]

        for part in parts:
            if not part: continue

            if part.startswith('[font'):
                new_style = style_stack[-1].copy()
                attrs = {}
                attr_string = part[part.find(' '):-1].strip()
                pairs = re.findall(r'(\w+)=(".*?"|\'.*?\'|[^\s\']*)', attr_string)
                for key, value in pairs:
                    attrs[key.lower()] = value.strip("'\"")

                if 'bold' in attrs: new_style['bold'] = attrs['bold'].lower() == 'true'
                if 'italic' in attrs: new_style['italic'] = attrs['italic'].lower() == 'true'
                if 'underline' in attrs: new_style['underline'] = attrs['underline'].lower() == 'true'
                if 'name' in attrs: new_style['name'] = attrs['name']
                if 'size' in attrs:
                    try: new_style['size'] = DocxPt(int(attrs['size']))  # type: ignore
                    except ValueError: pass
                if 'color' in attrs:
                    try: new_style['color'] = DocxRGBColor.from_string(attrs['color'].replace("#", ""))  # type: ignore
                    except (ValueError, KeyError): pass
                style_stack.append(new_style)

            elif part == '[/font]':
                if len(style_stack) > 1:
                    style_stack.pop()

            else:
                run = p.add_run()
                run.text = part
                current_style = style_stack[-1]
                font = run.font
                if current_style['bold'] is not None: font.bold = current_style['bold']
                if current_style['italic'] is not None: font.italic = current_style['italic']
                if current_style['underline'] is not None: font.underline = current_style['underline']
                if current_style['name'] is not None: font.name = current_style['name']
                if current_style['size'] is not None: font.size = current_style['size']  # type: ignore
                if current_style['color'] is not None: font.color.rgb = current_style['color']  # type: ignore

def split_content_and_notes(raw_content: str, title: str) -> tuple[list[str], list[str]]:
    """
    Splits raw slide content into main content lines and speaker notes.
    Lines after 'Speaker Notes:' or 'Notes:' are treated as notes.
    """
    content_lines = []
    notes_lines = []
    notes_section = False
    for line in raw_content.splitlines():
        stripped = line.strip()
        if stripped.lower().startswith('speaker notes:') or stripped.lower().startswith('notes:'):
            notes_section = True
            continue
        if notes_section:
            notes_lines.append(line)
        else:
            content_lines.append(line)
    return content_lines, notes_lines

# --- UI for Modes ---

def render_presentation_mode(dbms):
    st.header("‚ú® Presentation Generator")
    init_presentation_state()

    # STEP 1: Get topic and generate outline
    if st.session_state.presentation_step == "initial":
        st.subheader("Step 1: Choose a Topic")
        topic = st.text_input("What is your presentation about?", key="ppt_topic_input")
        if st.button("Generate Outline") and topic:
            st.session_state.presentation_topic = topic
            with st.spinner("Analyzing your documents and creating an outline..."):
                context = get_context_for_topic(dbms, topic)
                prompt = (f"Generate a slide-by-slide outline for a presentation on '{topic}'. "
                          f"Format it in markdown with each slide title starting with '#### ' and bullet points with '- '. "
                          f"Base the outline on the following information from my documents:\n\n{context}")
                
                # Using a placeholder for the system message
                outline_messages: list[ChatCompletionMessageParam] = [
                    {"role": "system", "content": "You are an AI assistant that creates presentation outlines."},
                    {"role": "user", "content": prompt}
                ]
                
                outline = "".join(list(openai_api_call(outline_messages, "Idx")))
                st.session_state.presentation_outline = outline
                st.session_state.presentation_step = "outline_generated"
                st.rerun()

    # STEP 2: Review and edit outline
    elif st.session_state.presentation_step == "outline_generated":
        st.subheader("Step 2: Edit Outline")
        outline = st.text_area("Edit the Markdown outline", value=st.session_state.presentation_outline, height=300)
        if st.button("Confirm Outline"):
            st.session_state.presentation_outline = outline
            parsed = parse_outline_to_slides(outline)
            # initialize content list and start page-by-page generation
            st.session_state.generated_content = []
            st.session_state.total_slides = len(parsed)
            st.session_state.parsed_slides = parsed
            st.session_state.slide_gen_index = 0
            st.session_state.presentation_step = "generating_slide_content"
            st.rerun()

    # STEP 3: Auto-generate all slide content
    elif st.session_state.presentation_step == "generating_slide_content":
        parsed = st.session_state.parsed_slides
        total = len(parsed)
        generated = []
        st.subheader(f"Step 3: Generating {total} slides...")
        progress = st.progress(0)
        with st.spinner("Generating slide content..."):
            for i, slide_info in enumerate(parsed):
                # Enhanced prompt: include formatting and centering instructions
                prompt = (
                    f"Write detailed speaker notes and bullet points for a slide titled '{slide_info['title']}'. "
                    f"Use markdown formatting: **bold**, *italic*, and bullet points with '- '. "
                    f"Include '[align=center]' tags to center the slide title. "
                    f"Include font size tags (e.g., '[font size=18]') where appropriate for slide display. "
                    f"Points: {', '.join(slide_info['points'])}."
                )
                messages = [
                    {"role": "system", "content": "You are an AI assistant that outputs slide content using custom formatting tags."},
                    {"role": "user", "content": prompt}
                ]
                content = "".join(openai_api_call(messages, "Idx"))
                generated.append({"title": slide_info['title'], "content": content})
                progress.progress(int((i+1)/total * 100))
        # Persist generated content and move to review
        st.session_state.presentation_content = generated
        st.session_state.presentation_step = "content_generated"
        # Clean up temporary state
        for key in ('generated_content','parsed_slides','total_slides','slide_gen_index'):
            st.session_state.pop(key, None)
        st.rerun()

    # STEP 4: Review & Edit Final Slide Content
    elif st.session_state.presentation_step == "content_generated":
        st.subheader("Step 4: Edit Slide Content")
        new_content = []
        for idx, slide in enumerate(st.session_state.presentation_content):
            st.markdown(f"**Slide {idx+1}: {slide['title']}**")
            content = st.text_area(f"Content for slide {idx+1}", value=slide['content'], key=f"slide_{idx}_content", height=200)
            new_content.append({"title": slide['title'], "content": content})
        if st.button("Generate Presentation"):
            st.session_state.presentation_content = new_content
            st.session_state.presentation_step = "export"
            st.rerun()

    # STEP 5: Export Presentation
    elif st.session_state.presentation_step == "export":
        st.subheader("Export Presentation")
        topic = st.session_state.presentation_topic
        col1, col2 = st.columns(2)
        with col1:
            ppt_path = create_presentation_from_content(st.session_state.presentation_content, topic, os.getenv("DASHSCOPE_API_KEY", ""))
            with open(ppt_path, "rb") as f:
                st.download_button("‚¨áÔ∏è Download PowerPoint", f, file_name=os.path.basename(ppt_path), mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        with col2:
            doc_path = create_document_from_content(st.session_state.presentation_content, topic)
            with open(doc_path, "rb") as f:
                st.download_button("‚¨áÔ∏è Download Word Document", f, file_name=os.path.basename(doc_path), mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

        if st.button("‚ú® Create Another Presentation"):
            init_presentation_state(force_reset=True)
            st.rerun()

def render_study_guide_mode(dbms):
    st.header(t("sg_header"))
    init_study_guide_state()

    # STEP 1: Get topic, style, and generate outline
    if st.session_state.study_guide_step == "initial":
        st.subheader(t("sg_configure"))
        
        col1, col2 = st.columns([2, 1])
        with col1:
            topic = st.text_input(t("sg_topic_input"), key="study_guide_topic_input", 
                                placeholder=t("sg_topic_placeholder"))
        
        with col2:
            style = st.selectbox(t("sg_style"), 
                               ["comprehensive", "summary", "outline", "flashcard_prep", "exam_focused"],
                               format_func=lambda x: x.replace('_', ' ').title(),
                               key="study_guide_style_select")
        
        # Style descriptions
        style_descriptions = {
            "comprehensive": "üìñ Detailed explanations with examples and practice questions",
            "summary": "üìù Concise key points and essential information",
            "outline": "üìã Structured hierarchical organization of topics", 
            "flashcard_prep": "üÉè Optimized for creating flashcards with clear Q&A format",
            "exam_focused": "üéØ Emphasis on likely test questions and critical concepts"
        }
        st.info(style_descriptions[style])
        
        # Extreme Mode checkbox
        st.markdown("### ‚öôÔ∏è Generation Options")
        col1, col2 = st.columns(2)
        
        with col1:
            extreme_mode = st.checkbox(
                "üöÄ Extreme Mode", 
                value=False,
                key="extreme_mode_checkbox",
                help="‚ö†Ô∏è Extreme Mode uses advanced AI processing for maximum detail and quality but may take up to 10 minutes to generate. Unchecked uses fast mode (IDX) which generates content in under 2 minutes."
            )
        
        with col2:
            page_count = st.number_input(
                "üìÑ Target Pages", 
                min_value=-1, 
                max_value=50, 
                value=-1,
                step=1,
                key="page_count_input",
                help="‚ùì Set target page count for the study guide. Use -1 for automatic length (AI decides optimal length). Range: 1-50 pages. Higher page counts will take longer to generate."
            )
        
        # Show estimates based on page count
        if page_count == -1:
            st.info("üìù **Auto Length**: AI will determine optimal study guide length based on topic complexity.")
        elif page_count <= 5:
            st.success(f"üìÑ **Short Guide**: ~{page_count} pages ‚Ä¢ Quick generation ‚Ä¢ Concise content")
        elif page_count <= 15:
            st.info(f"üìÑ **Medium Guide**: ~{page_count} pages ‚Ä¢ Moderate generation time ‚Ä¢ Balanced detail")
        elif page_count <= 30:
            st.warning(f"üìÑ **Long Guide**: ~{page_count} pages ‚Ä¢ Extended generation time ‚Ä¢ Comprehensive detail")
        else:
            st.error(f"üìÑ **Very Long Guide**: ~{page_count} pages ‚Ä¢ Very long generation time ‚Ä¢ Extensive detail")
        
        # Show time estimates
        if extreme_mode:
            if page_count == -1:
                time_est = "5-15 minutes"
            elif page_count <= 10:
                time_est = f"{3 + page_count//2}-{5 + page_count} minutes"
            else:
                time_est = f"{5 + page_count//2}-{10 + page_count} minutes"
            st.warning(f"‚ö†Ô∏è **Extreme Mode Enabled**: Generation may take {time_est} but will produce the highest quality, most comprehensive content.")
        else:
            if page_count == -1:
                time_est = "1-3 minutes"
            elif page_count <= 15:
                time_est = f"{1 + page_count//5}-{2 + page_count//3} minutes"
            else:
                time_est = f"{3 + page_count//5}-{5 + page_count//3} minutes"
            st.success(f"‚ö° **Fast Mode**: Using optimized IDX processing for quick generation (Est. {time_est}).")
        
        if st.button("Generate Study Guide Outline") and topic:
            st.session_state.study_guide_topic = topic
            st.session_state.study_guide_style = style
            st.session_state.study_guide_extreme_mode = extreme_mode
            st.session_state.study_guide_page_count = page_count
            
            generation_mode = "Normal" if extreme_mode else "Idx"
            time_estimate = "5-10 minutes" if extreme_mode else "under 45 seconds per page"
            
            with st.spinner(f"Analyzing your documents and creating a study guide outline... (Estimated time: {time_estimate})"):
                context = get_context_for_topic(dbms, topic)
                style_instruction = get_study_guide_style_prompt(style)
                
                # Add page count guidance to outline generation
                if page_count == -1:
                    length_guidance = "Create an outline with natural section count based on topic complexity."
                else:
                    recommended_sections = max(3, min(page_count // 2, 12))  # 2-12 sections based on page count
                    length_guidance = f"Create an outline with approximately {recommended_sections} main sections to fit ~{page_count} pages."
                
                prompt = (f"Create a detailed outline for a study guide on '{topic}'. "
                         f"{style_instruction} "
                         f"{length_guidance} "
                         f"Format the outline in markdown with main sections starting with '## ' and subtopics with '- '. "
                         f"Use proper markdown formatting throughout. "
                         f"Base the outline on the following information from my documents:\n\n{context}")
                
                outline_messages: list[ChatCompletionMessageParam] = [
                    {"role": "system", "content": "You are an educational AI that creates comprehensive study guide outlines."},
                    {"role": "user", "content": prompt}
                ]
                
                outline = "".join(list(openai_api_call(outline_messages, generation_mode)))
                st.session_state.study_guide_outline = outline
                st.session_state.study_guide_step = "outline_generated"
                st.rerun()

    # STEP 2: Review and edit outline
    elif st.session_state.study_guide_step == "outline_generated":
        st.subheader("Step 2: Review & Edit Outline")
        
        # Show current settings
        col1, col2, col3 = st.columns(3)
        with col1:
            st.info(f"üìö Style: {st.session_state.study_guide_style.replace('_', ' ').title()}")
        with col2:
            mode_display = "üöÄ Extreme Mode" if st.session_state.study_guide_extreme_mode else "‚ö° Fast Mode"
            st.info(f"‚öôÔ∏è {mode_display}")
        with col3:
            page_count = st.session_state.get('study_guide_page_count', -1)
            page_display = "Auto Length" if page_count == -1 else f"{page_count} Pages"
            st.info(f"üìÑ Target: {page_display}")
        
        outline = st.text_area("Edit the study guide outline:", 
                             value=st.session_state.study_guide_outline, 
                             height=400,
                             help="Use ## for main sections and - for subtopics")
        
        col1, col2 = st.columns(2)
        with col1:
            if st.button("‚¨ÖÔ∏è Back to Topic Selection"):
                st.session_state.study_guide_step = "initial"
                st.rerun()
        
        with col2:
            if st.button("Generate Content ‚û°Ô∏è"):
                st.session_state.study_guide_outline = outline
                parsed = parse_outline_to_sections(outline)
                st.session_state.parsed_sections = parsed
                st.session_state.study_guide_step = "generating_content"
                st.rerun()

    # STEP 3: Auto-generate all section content
    elif st.session_state.study_guide_step == "generating_content":
        parsed = st.session_state.parsed_sections
        total = len(parsed)
        
        # Determine generation mode and estimates
        is_extreme_mode = st.session_state.get('study_guide_extreme_mode', False)
        generation_mode = "Normal" if is_extreme_mode else "Idx"
        mode_display = "üöÄ Extreme Mode" if is_extreme_mode else "‚ö° Fast Mode"
        time_estimate = "5-10 minutes" if is_extreme_mode else "under 2 minutes"
        
        st.subheader(f"Step 3: Generating {total} sections...")
        
        # Show generation info
        col1, col2, col3 = st.columns(3)
        with col1:
            st.info(f"üìö {st.session_state.study_guide_style.replace('_', ' ').title()} style")
        with col2:
            st.info(f"‚öôÔ∏è {mode_display} ‚Ä¢ Est. {time_estimate}")
        with col3:
            page_count = st.session_state.get('study_guide_page_count', -1)
            if page_count == -1:
                st.info("üìÑ Auto Length")
            else:
                st.info(f"üìÑ ~{page_count} pages")
        
        if is_extreme_mode:
            st.warning("üöÄ **Extreme Mode Active**: Using advanced AI processing for maximum detail. Please be patient as this may take several minutes.")
        else:
            st.success("‚ö° **Fast Mode Active**: Using optimized processing for quick generation.")
        
        progress = st.progress(0)
        generated = []
        
        with st.spinner(f"Generating detailed study guide content using {mode_display}..."):
            context = get_context_for_topic(dbms, st.session_state.study_guide_topic)
            style_instruction = get_study_guide_style_prompt(st.session_state.study_guide_style)
            page_count = st.session_state.get('study_guide_page_count', -1)
            page_count_instruction = get_page_count_instructions(page_count, len(parsed))
            
            for i, section_info in enumerate(parsed):
                # Enhanced prompt for study guide content
                markdown_instructions = get_markdown_formatting_instructions()
                prompt = (
                    f"Write comprehensive study guide content for the section titled '{section_info['title']}'. "
                    f"{style_instruction} "
                    f"Include the following subtopics: {', '.join(section_info['subsections'])}. "
                    f"{markdown_instructions} "
                    f"{page_count_instruction} "
                    f"Include definitions, explanations, examples, and key points. "
                    f"Add practice questions or review items where appropriate. "
                    f"Structure your response with clear headers and organized sections. "
                    f"Base your content on this context: {context[:2000]}..."  # Limit context to avoid token limits
                )
                
                messages = [
                    {"role": "system", "content": "You are an expert educator creating detailed study guide content."},
                    {"role": "user", "content": prompt}
                ]
                
                content = "".join(openai_api_call(messages, generation_mode))
                generated.append({"title": section_info['title'], "content": content})
                progress.progress(int((i+1)/total * 100))
        
        # Store generated content and move to review step
        st.session_state.study_guide_content = generated
        st.session_state.study_guide_step = "content_generated"
        
        # Clean up temporary state
        st.session_state.pop('parsed_sections', None)
        st.rerun()

    # STEP 4: Review & Edit Final Content
    elif st.session_state.study_guide_step == "content_generated":
        st.subheader("Step 4: Review & Edit Content")
        
        # Show generation info
        col1, col2, col3 = st.columns(3)
        with col1:
            st.info(f"üìö Topic: {st.session_state.study_guide_topic}")
        with col2:
            style_display = st.session_state.study_guide_style.replace('_', ' ').title()
            mode_display = "üöÄ Extreme Mode" if st.session_state.get('study_guide_extreme_mode', False) else "‚ö° Fast Mode"
            st.info(f"‚öôÔ∏è {style_display} ‚Ä¢ {mode_display}")
        with col3:
            page_count = st.session_state.get('study_guide_page_count', -1)
            if page_count == -1:
                st.info("üìÑ Auto Length")
            else:
                st.info(f"üìÑ ~{page_count} pages")
        
        # Add tabs for easier navigation through sections
        if len(st.session_state.study_guide_content) > 1:
            section_tabs = st.tabs([f"Section {i+1}: {section['title'][:20]}..." if len(section['title']) > 20 
                                  else f"Section {i+1}: {section['title']}" 
                                  for i, section in enumerate(st.session_state.study_guide_content)])
            
            new_content = []
            for idx, (tab, section) in enumerate(zip(section_tabs, st.session_state.study_guide_content)):
                with tab:
                    st.markdown(f"**{section['title']}**")
                    content = st.text_area(
                        f"Content for {section['title']}", 
                        value=section['content'], 
                        key=f"section_{idx}_content", 
                        height=300
                    )
                    new_content.append({"title": section['title'], "content": content})
        else:
            # Single section - no tabs needed
            section = st.session_state.study_guide_content[0]
            st.markdown(f"**{section['title']}**")
            content = st.text_area(
                f"Content for {section['title']}", 
                value=section['content'], 
                key="section_0_content", 
                height=400
            )
            new_content = [{"title": section['title'], "content": content}]
        
        col1, col2 = st.columns(2)
        with col1:
            if st.button("‚¨ÖÔ∏è Back to Outline"):
                st.session_state.study_guide_step = "outline_generated"
                st.rerun()
        
        with col2:
            if st.button("Generate Study Guide üìö"):
                st.session_state.study_guide_content = new_content
                st.session_state.study_guide_step = "export"
                st.rerun()

    # STEP 5: Export Study Guide
    elif st.session_state.study_guide_step == "export":
        st.subheader("üì• Download Your Study Guide")
        
        topic = st.session_state.study_guide_topic
        style = st.session_state.study_guide_style
        mode_display = "üöÄ Extreme Mode" if st.session_state.get('study_guide_extreme_mode', False) else "‚ö° Fast Mode"
        
        # Show final summary
        st.success(f"‚úÖ **Study Guide Complete!** Generated using {mode_display}")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("### üìÑ Word Document")
            st.write("Complete study guide with formatting and table of contents")
            
            try:
                doc_path = create_study_guide_from_content(
                    st.session_state.study_guide_content, 
                    topic, 
                    style,
                    st.session_state.get('study_guide_extreme_mode', False)
                )
                with open(doc_path, "rb") as f:
                    st.download_button(
                        "‚¨áÔ∏è Download Study Guide (.docx)", 
                        f, 
                        file_name=os.path.basename(doc_path), 
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        use_container_width=True
                    )
            except Exception as e:
                st.error(f"Error creating document: {str(e)}")
        
        with col2:
            st.markdown("### üìã Text Format")
            st.write("Plain text version for easy copying")
            
            # Create text version with generation mode info
            text_content = f"STUDY GUIDE: {topic.upper()}\n"
            text_content += f"Generated: {datetime.datetime.now().strftime('%B %d, %Y')}\n"
            text_content += f"Style: {style.replace('_', ' ').title()}\n"
            text_content += f"Generation Mode: {mode_display}\n"
            text_content += "=" * 50 + "\n\n"
            
            for section in st.session_state.study_guide_content:
                text_content += f"{section['title']}\n"
                text_content += "-" * len(section['title']) + "\n"
                # Remove formatting tags for plain text
                clean_content = re.sub(r'\[/?(font|align)[^]]*\]', '', section['content'])
                text_content += clean_content + "\n\n"
            
            st.download_button(
                "‚¨áÔ∏è Download as Text (.txt)",
                text_content.encode('utf-8'),
                file_name=f"StudyGuide_{topic.replace(' ', '_')}_{style}.txt",
                mime="text/plain",
                use_container_width=True
            )
        
        st.markdown("---")
        col1, col2 = st.columns(2)
        with col1:
            if st.button("üìö Create Another Study Guide"):
                init_study_guide_state(force_reset=True)
                st.rerun()
        
        with col2:
            if st.button("üìù Create Presentation from This Topic"):
                # Switch to presentation mode with the same topic
                init_presentation_state(force_reset=True)
                st.session_state.presentation_topic = topic
                st.session_state.mixup_mode_selector = "Presentation"  # Switch the radio button
                st.rerun()

def render_flashcard_mode():
    st.header("üìá Q&A Flashcard Generator")
    st.info("‚ö†Ô∏è **Note:** This is a simplified implementation. Full flashcard functionality coming soon!")
    
    # Simple flashcard generator
    st.subheader("Quick Flashcard Generator")
    
    col1, col2 = st.columns([2, 1])
    with col1:
        topic = st.text_input("Enter a topic for flashcard generation:", 
                            placeholder="e.g., Biology Terms, History Dates, Math Formulas")
    
    with col2:
        num_cards = st.selectbox("Number of flashcards:", [5, 10, 15, 20], index=1)
    
    if st.button("Generate Flashcards") and topic:
        with st.spinner(f"Generating {num_cards} flashcards for {topic}..."):
            # Simple prompt for flashcard generation
            prompt = (f"Create {num_cards} educational flashcards about {topic}. "
                     f"Format each flashcard as 'Q: [question]' followed by 'A: [answer]' on the next line. "
                     f"Make questions clear and concise, and answers accurate and helpful for studying.")
            
            messages: list[ChatCompletionMessageParam] = [
                {"role": "system", "content": "You are an educational AI that creates study flashcards."},
                {"role": "user", "content": prompt}
            ]
            
            flashcard_content = "".join(openai_api_call(messages, "Idx"))
            
            # Parse and display flashcards
            lines = flashcard_content.split('\n')
            current_question = ""
            current_answer = ""
            flashcards = []
            
            for line in lines:
                line = line.strip()
                if line.startswith('Q:'):
                    if current_question and current_answer:
                        flashcards.append({'question': current_question, 'answer': current_answer})
                    current_question = line[2:].strip()
                    current_answer = ""
                elif line.startswith('A:'):
                    current_answer = line[2:].strip()
                elif current_answer and line:  # Continue answer on next lines
                    current_answer += " " + line
            
            # Add the last flashcard
            if current_question and current_answer:
                flashcards.append({'question': current_question, 'answer': current_answer})
            
            # Display flashcards
            if flashcards:
                st.success(f"Generated {len(flashcards)} flashcards!")
                
                # Create tabs for each flashcard
                if len(flashcards) > 1:
                    flashcard_tabs = st.tabs([f"Card {i+1}" for i in range(len(flashcards))])
                    
                    for i, (tab, card) in enumerate(zip(flashcard_tabs, flashcards)):
                        with tab:
                            st.markdown(f"### üÉè Flashcard {i+1}")
                            
                            # Create expandable sections for question and answer
                            with st.expander("‚ùì Question", expanded=True):
                                st.write(card['question'])
                            
                            with st.expander("‚úÖ Answer", expanded=False):
                                st.write(card['answer'])
                else:
                    # Single flashcard
                    card = flashcards[0]
                    st.markdown("### üÉè Flashcard")
                    
                    with st.expander("‚ùì Question", expanded=True):
                        st.write(card['question'])
                    
                    with st.expander("‚úÖ Answer", expanded=False):
                        st.write(card['answer'])
                
                # Export options
                st.markdown("---")
                st.subheader("üì• Export Flashcards")
                
                # Create text format for download
                flashcard_text = f"FLASHCARDS: {topic.upper()}\n"
                flashcard_text += f"Generated: {datetime.datetime.now().strftime('%B %d, %Y')}\n"
                flashcard_text += "=" * 50 + "\n\n"
                
                for i, card in enumerate(flashcards, 1):
                    flashcard_text += f"CARD {i}\n"
                    flashcard_text += f"Q: {card['question']}\n"
                    flashcard_text += f"A: {card['answer']}\n\n"
                
                st.download_button(
                    "‚¨áÔ∏è Download Flashcards (.txt)",
                    flashcard_text.encode('utf-8'),
                    file_name=f"Flashcards_{topic.replace(' ', '_')}.txt",
                    mime="text/plain"
                )
            else:
                st.error("Could not generate flashcards. Please try again with a different topic.")
    
    st.markdown("---")
    st.info("üí° **Tip:** For more advanced flashcard features, consider using the Study Guide mode and then creating flashcards from the generated content!")

def mixup_page():
    # Initialize language state and render language selector
    init_language_state()
    render_language_selector()
    
    st.title(t("title"))
    st.write(t("subtitle"))

    # Check if the database is uninitialized or has no indexed entries.
    if 'dbms' not in st.session_state or not isinstance(st.session_state.dbms, FiberDBMS) or st.session_state.dbms.is_empty():
        st.info(t("no_files_indexed"))
        # If the dbms object doesn't exist at all, create an empty one to prevent errors.
        if 'dbms' not in st.session_state or not isinstance(st.session_state.dbms, FiberDBMS):
            st.session_state.dbms = FiberDBMS()

    dbms = st.session_state.dbms

    mode = st.radio(
        t("choose_mode"),
        [t("presentation"), t("study_guide"), t("flashcards")],
        key='mixup_mode_selector',
        horizontal=True,
    )
    
    st.markdown("---")

    if mode == t("presentation"):
        render_presentation_mode(dbms)
    elif mode == t("study_guide"):
        render_study_guide_mode(dbms)
    elif mode == t("flashcards"):
        render_flashcard_mode()

# --- Translation System for UI ---

# Supported languages and their translations
TRANSLATIONS = {
    "en": {
        # Main titles and navigation
        "title": "Arcana Mixup",
        "subtitle": "Your intelligent assistant for creating documents, presentations, and study materials.",
        "choose_mode": "Choose a generation mode:",
        "language_selector": "üåê Interface Language:",
        
        # Mode names
        "presentation": "Presentation",
        "study_guide": "Study Guide", 
        "flashcards": "Q&A Flashcards",
        
        # Common buttons
        "back": "‚¨ÖÔ∏è Back",
        "next": "‚û°Ô∏è Next",
        "generate": "Generate",
        "download": "‚¨áÔ∏è Download",
        "create_another": "Create Another",
        "confirm": "Confirm",
        "edit": "Edit",
        
        # Study Guide specific
        "sg_header": "üìö Study Guide Generator",
        "sg_configure": "Step 1: Configure Your Study Guide",
        "sg_topic_input": "What topic would you like to study?",
        "sg_topic_placeholder": "e.g., Cell Biology, World War II, Calculus Derivatives",
        "sg_style": "Study Guide Style:",
        "sg_generation_options": "‚öôÔ∏è Generation Options",
        "sg_extreme_mode": "üöÄ Extreme Mode",
        "sg_target_pages": "üìÑ Target Pages",
        "sg_auto_length": "üìù **Auto Length**: AI will determine optimal study guide length based on topic complexity.",
        "sg_short_guide": "üìÑ **Short Guide**: ~{} pages ‚Ä¢ Quick generation ‚Ä¢ Concise content",
        "sg_medium_guide": "üìÑ **Medium Guide**: ~{} pages ‚Ä¢ Moderate generation time ‚Ä¢ Balanced detail", 
        "sg_long_guide": "üìÑ **Long Guide**: ~{} pages ‚Ä¢ Extended generation time ‚Ä¢ Comprehensive detail",
        "sg_very_long_guide": "üìÑ **Very Long Guide**: ~{} pages ‚Ä¢ Very long generation time ‚Ä¢ Extensive detail",
        
        # Presentation specific
        "ppt_header": "‚ú® Presentation Generator",
        "ppt_choose_topic": "Step 1: Choose a Topic",
        "ppt_topic_input": "What is your presentation about?",
        
        # Common messages
        "no_files_indexed": "No indexed files found. Content will be generated from general knowledge. To use your own documents as context, please go to the 'Files' page and index them first.",
        "generating": "Generating...",
        "analyzing": "Analyzing your documents...",
        "complete": "‚úÖ Complete!",
        "error": "‚ùå Error:",
        
        # File types
        "word_document": "üìÑ Word Document",
        "powerpoint": "üìä PowerPoint",
        "text_format": "üìã Text Format",
    },
    
    "es": {
        # Main titles and navigation
        "title": "Arcana Mixup",
        "subtitle": "Tu asistente inteligente para crear documentos, presentaciones y materiales de estudio.",
        "choose_mode": "Elige un modo de generaci√≥n:",
        "language_selector": "üåê Idioma de Interfaz:",
        
        # Mode names
        "presentation": "Presentaci√≥n",
        "study_guide": "Gu√≠a de Estudio",
        "flashcards": "Tarjetas P&R",
        
        # Common buttons
        "back": "‚¨ÖÔ∏è Atr√°s",
        "next": "‚û°Ô∏è Siguiente", 
        "generate": "Generar",
        "download": "‚¨áÔ∏è Descargar",
        "create_another": "Crear Otro",
        "confirm": "Confirmar",
        "edit": "Editar",
        
        # Study Guide specific
        "sg_header": "üìö Generador de Gu√≠as de Estudio",
        "sg_configure": "Paso 1: Configura tu Gu√≠a de Estudio",
        "sg_topic_input": "¬øQu√© tema te gustar√≠a estudiar?",
        "sg_topic_placeholder": "ej., Biolog√≠a Celular, Segunda Guerra Mundial, Derivadas de C√°lculo",
        "sg_style": "Estilo de Gu√≠a de Estudio:",
        "sg_generation_options": "‚öôÔ∏è Opciones de Generaci√≥n",
        "sg_extreme_mode": "üöÄ Modo Extremo",
        "sg_target_pages": "üìÑ P√°ginas Objetivo",
        "sg_auto_length": "üìù **Longitud Autom√°tica**: La IA determinar√° la longitud √≥ptima basada en la complejidad del tema.",
        "sg_short_guide": "üìÑ **Gu√≠a Corta**: ~{} p√°ginas ‚Ä¢ Generaci√≥n r√°pida ‚Ä¢ Contenido conciso",
        "sg_medium_guide": "üìÑ **Gu√≠a Media**: ~{} p√°ginas ‚Ä¢ Tiempo moderado ‚Ä¢ Detalle equilibrado",
        "sg_long_guide": "üìÑ **Gu√≠a Larga**: ~{} p√°ginas ‚Ä¢ Tiempo extendido ‚Ä¢ Detalle comprensivo",
        "sg_very_long_guide": "üìÑ **Gu√≠a Muy Larga**: ~{} p√°ginas ‚Ä¢ Tiempo muy largo ‚Ä¢ Detalle extensivo",
        
        # Presentation specific
        "ppt_header": "‚ú® Generador de Presentaciones",
        "ppt_choose_topic": "Paso 1: Elige un Tema",
        "ppt_topic_input": "¬øDe qu√© trata tu presentaci√≥n?",
        
        # Common messages
        "no_files_indexed": "No se encontraron archivos indexados. El contenido se generar√° desde conocimiento general. Para usar tus propios documentos como contexto, ve a la p√°gina 'Archivos' e indexa primero.",
        "generating": "Generando...",
        "analyzing": "Analizando tus documentos...",
        "complete": "‚úÖ ¬°Completo!",
        "error": "‚ùå Error:",
        
        # File types
        "word_document": "üìÑ Documento Word",
        "powerpoint": "üìä PowerPoint", 
        "text_format": "üìã Formato Texto",
    },
    
    "fr": {
        # Main titles and navigation
        "title": "Arcana Mixup",
        "subtitle": "Votre assistant intelligent pour cr√©er des documents, pr√©sentations et mat√©riels d'√©tude.",
        "choose_mode": "Choisissez un mode de g√©n√©ration:",
        "language_selector": "üåê Langue d'Interface:",
        
        # Mode names
        "presentation": "Pr√©sentation",
        "study_guide": "Guide d'√âtude",
        "flashcards": "Cartes Q&R",
        
        # Common buttons
        "back": "‚¨ÖÔ∏è Retour",
        "next": "‚û°Ô∏è Suivant",
        "generate": "G√©n√©rer",
        "download": "‚¨áÔ∏è T√©l√©charger",
        "create_another": "Cr√©er un Autre",
        "confirm": "Confirmer",
        "edit": "Modifier",
        
        # Study Guide specific
        "sg_header": "üìö G√©n√©rateur de Guides d'√âtude",
        "sg_configure": "√âtape 1: Configurez votre Guide d'√âtude",
        "sg_topic_input": "Quel sujet aimeriez-vous √©tudier?",
        "sg_topic_placeholder": "ex., Biologie Cellulaire, Seconde Guerre Mondiale, D√©riv√©es du Calcul",
        "sg_style": "Style de Guide d'√âtude:",
        "sg_generation_options": "‚öôÔ∏è Options de G√©n√©ration",
        "sg_extreme_mode": "üöÄ Mode Extr√™me",
        "sg_target_pages": "üìÑ Pages Cibles",
        "sg_auto_length": "üìù **Longueur Automatique**: L'IA d√©terminera la longueur optimale bas√©e sur la complexit√© du sujet.",
        "sg_short_guide": "üìÑ **Guide Court**: ~{} pages ‚Ä¢ G√©n√©ration rapide ‚Ä¢ Contenu concis",
        "sg_medium_guide": "üìÑ **Guide Moyen**: ~{} pages ‚Ä¢ Temps mod√©r√© ‚Ä¢ D√©tail √©quilibr√©",
        "sg_long_guide": "üìÑ **Guide Long**: ~{} pages ‚Ä¢ Temps √©tendu ‚Ä¢ D√©tail compr√©hensif",
        "sg_very_long_guide": "üìÑ **Guide Tr√®s Long**: ~{} pages ‚Ä¢ Temps tr√®s long ‚Ä¢ D√©tail extensif",
        
        # Presentation specific
        "ppt_header": "‚ú® G√©n√©rateur de Pr√©sentations",
        "ppt_choose_topic": "√âtape 1: Choisissez un Sujet",
        "ppt_topic_input": "De quoi parle votre pr√©sentation?",
        
        # Common messages
        "no_files_indexed": "Aucun fichier index√© trouv√©. Le contenu sera g√©n√©r√© √† partir de connaissances g√©n√©rales. Pour utiliser vos propres documents comme contexte, allez √† la page 'Fichiers' et indexez d'abord.",
        "generating": "G√©n√©ration...",
        "analyzing": "Analyse de vos documents...",
        "complete": "‚úÖ Termin√©!",
        "error": "‚ùå Erreur:",
        
        # File types
        "word_document": "üìÑ Document Word",
        "powerpoint": "üìä PowerPoint",
        "text_format": "üìã Format Texte",
    },
    
    "de": {
        # Main titles and navigation
        "title": "Arcana Mixup",
        "subtitle": "Ihr intelligenter Assistent f√ºr das Erstellen von Dokumenten, Pr√§sentationen und Lernmaterialien.",
        "choose_mode": "W√§hlen Sie einen Generierungsmodus:",
        "language_selector": "üåê Interface-Sprache:",
        
        # Mode names
        "presentation": "Pr√§sentation",
        "study_guide": "Lernleitfaden",
        "flashcards": "F&A Karten",
        
        # Common buttons
        "back": "‚¨ÖÔ∏è Zur√ºck",
        "next": "‚û°Ô∏è Weiter",
        "generate": "Generieren",
        "download": "‚¨áÔ∏è Herunterladen",
        "create_another": "Weitere Erstellen",
        "confirm": "Best√§tigen",
        "edit": "Bearbeiten",
        
        # Study Guide specific
        "sg_header": "üìö Lernleitfaden-Generator",
        "sg_configure": "Schritt 1: Konfigurieren Sie Ihren Lernleitfaden",
        "sg_topic_input": "Welches Thema m√∂chten Sie studieren?",
        "sg_topic_placeholder": "z.B., Zellbiologie, Zweiter Weltkrieg, Kalk√ºl-Ableitungen",
        "sg_style": "Lernleitfaden-Stil:",
        "sg_generation_options": "‚öôÔ∏è Generierungsoptionen",
        "sg_extreme_mode": "üöÄ Extremmodus",
        "sg_target_pages": "üìÑ Zielseiten",
        "sg_auto_length": "üìù **Automatische L√§nge**: KI bestimmt optimale Leitfadenl√§nge basierend auf Themenkomplexit√§t.",
        "sg_short_guide": "üìÑ **Kurzer Leitfaden**: ~{} Seiten ‚Ä¢ Schnelle Generierung ‚Ä¢ Pr√§gnanter Inhalt",
        "sg_medium_guide": "üìÑ **Mittlerer Leitfaden**: ~{} Seiten ‚Ä¢ Moderate Zeit ‚Ä¢ Ausgewogenes Detail",
        "sg_long_guide": "üìÑ **Langer Leitfaden**: ~{} Seiten ‚Ä¢ Erweiterte Zeit ‚Ä¢ Umfassendes Detail",
        "sg_very_long_guide": "üìÑ **Sehr Langer Leitfaden**: ~{} Seiten ‚Ä¢ Sehr lange Zeit ‚Ä¢ Extensives Detail",
        
        # Presentation specific
        "ppt_header": "‚ú® Pr√§sentations-Generator",
        "ppt_choose_topic": "Schritt 1: W√§hlen Sie ein Thema",
        "ppt_topic_input": "Worum geht es in Ihrer Pr√§sentation?",
        
        # Common messages
        "no_files_indexed": "Keine indexierten Dateien gefunden. Inhalt wird aus allgemeinem Wissen generiert. Um Ihre eigenen Dokumente als Kontext zu verwenden, gehen Sie zur 'Dateien'-Seite und indexieren Sie zuerst.",
        "generating": "Generierung...",
        "analyzing": "Analysiere Ihre Dokumente...",
        "complete": "‚úÖ Fertig!",
        "error": "‚ùå Fehler:",
        
        # File types
        "word_document": "üìÑ Word-Dokument",
        "powerpoint": "üìä PowerPoint",
        "text_format": "üìã Textformat",
    },
    
    "zh": {
        # Main titles and navigation
        "title": "ÈòøÂç°Á∫≥ÔºöÈöèÁ¨îÁ≥ªÁªü",
        "subtitle": "ÊÇ®ÂàõÂª∫ÊñáÊ°£„ÄÅÊºîÁ§∫ÊñáÁ®øÂíåÂ≠¶‰π†ÊùêÊñôÁöÑÊô∫ËÉΩÂä©Êâã„ÄÇ",
        "choose_mode": "ÈÄâÊã©ÁîüÊàêÊ®°ÂºèÔºö",
        "language_selector": "üåê ÁïåÈù¢ËØ≠Ë®ÄÔºö",
        
        # Mode names
        "presentation": "ÊºîÁ§∫ÊñáÁ®ø",
        "study_guide": "Â≠¶‰π†ÊåáÂçó",
        "flashcards": "ÈóÆÁ≠îÂç°Áâá",
        
        # Common buttons
        "back": "‚¨ÖÔ∏è ËøîÂõû",
        "next": "‚û°Ô∏è ‰∏ã‰∏ÄÊ≠•",
        "generate": "ÁîüÊàê",
        "download": "‚¨áÔ∏è ‰∏ãËΩΩ",
        "create_another": "ÂàõÂª∫Âè¶‰∏Ä‰∏™",
        "confirm": "Á°ÆËÆ§",
        "edit": "ÁºñËæë",
        
        # Study Guide specific
        "sg_header": "üìö Â≠¶‰π†ÊåáÂçóÁîüÊàêÂô®",
        "sg_configure": "Ê≠•È™§1ÔºöÈÖçÁΩÆÊÇ®ÁöÑÂ≠¶‰π†ÊåáÂçó",
        "sg_topic_input": "ÊÇ®ÊÉ≥Â≠¶‰π†‰ªÄ‰πà‰∏ªÈ¢òÔºü",
        "sg_topic_placeholder": "‰æãÂ¶ÇÔºöÁªÜËÉûÁîüÁâ©Â≠¶„ÄÅÁ¨¨‰∫åÊ¨°‰∏ñÁïåÂ§ßÊàò„ÄÅÂæÆÁßØÂàÜÂØºÊï∞",
        "sg_style": "Â≠¶‰π†ÊåáÂçóÈ£éÊ†ºÔºö",
        "sg_generation_options": "‚öôÔ∏è ÁîüÊàêÈÄâÈ°π",
        "sg_extreme_mode": "üöÄ ÊûÅÈôêÊ®°Âºè",
        "sg_target_pages": "üìÑ ÁõÆÊ†áÈ°µÊï∞",
        "sg_auto_length": "üìù **Ëá™Âä®ÈïøÂ∫¶**ÔºöAIÂ∞ÜÊ†πÊçÆ‰∏ªÈ¢òÂ§çÊùÇÊÄßÁ°ÆÂÆöÊúÄ‰Ω≥ÊåáÂçóÈïøÂ∫¶„ÄÇ",
        "sg_short_guide": "üìÑ **ÁÆÄÁü≠ÊåáÂçó**Ôºö~{}È°µ ‚Ä¢ Âø´ÈÄüÁîüÊàê ‚Ä¢ ÁÆÄÊ¥ÅÂÜÖÂÆπ",
        "sg_medium_guide": "üìÑ **‰∏≠Á≠âÊåáÂçó**Ôºö~{}È°µ ‚Ä¢ ÈÄÇ‰∏≠Êó∂Èó¥ ‚Ä¢ Âπ≥Ë°°ËØ¶ÁªÜ",
        "sg_long_guide": "üìÑ **ÈïøÊåáÂçó**Ôºö~{}È°µ ‚Ä¢ Êâ©Â±ïÊó∂Èó¥ ‚Ä¢ ÂÖ®Èù¢ËØ¶ÁªÜ",
        "sg_very_long_guide": "üìÑ **Ë∂ÖÈïøÊåáÂçó**Ôºö~{}È°µ ‚Ä¢ ÂæàÈïøÊó∂Èó¥ ‚Ä¢ ÂπøÊ≥õËØ¶ÁªÜ",
        
        # Presentation specific
        "ppt_header": "‚ú® ÊºîÁ§∫ÊñáÁ®øÁîüÊàêÂô®",
        "ppt_choose_topic": "Ê≠•È™§1ÔºöÈÄâÊã©‰∏ªÈ¢ò",
        "ppt_topic_input": "ÊÇ®ÁöÑÊºîÁ§∫ÊñáÁ®øÊòØÂÖ≥‰∫é‰ªÄ‰πàÁöÑÔºü",
        
        # Common messages
        "no_files_indexed": "Êú™ÊâæÂà∞Á¥¢ÂºïÊñá‰ª∂„ÄÇÂ∞Ü‰ªé‰∏ÄËà¨Áü•ËØÜÁîüÊàêÂÜÖÂÆπ„ÄÇË¶Å‰ΩøÁî®ÊÇ®Ëá™Â∑±ÁöÑÊñáÊ°£‰Ωú‰∏∫‰∏ä‰∏ãÊñáÔºåËØ∑ÂÖàËΩ¨Âà∞Êñá‰ª∂È°µÈù¢Âπ∂Âª∫Á´ãÁ¥¢Âºï„ÄÇ",
        "generating": "ÁîüÊàê‰∏≠...",
        "analyzing": "ÂàÜÊûêÊÇ®ÁöÑÊñáÊ°£...",
        "complete": "‚úÖ ÂÆåÊàêÔºÅ",
        "error": "‚ùå ÈîôËØØÔºö",
        
        # File types
        "word_document": "üìÑ WordÊñáÊ°£",
        "powerpoint": "üìä PowerPoint",
        "text_format": "üìã ÊñáÊú¨Ê†ºÂºè",
    }
}

def get_available_languages():
    """Returns list of available language codes and names."""
    return {
        "en": "English",
        "es": "Espa√±ol", 
        "fr": "Fran√ßais",
        "de": "Deutsch",
        "zh": "‰∏≠Êñá"
    }

def init_language_state():
    """Initialize language selection in session state."""
    if 'ui_language' not in st.session_state:
        st.session_state.ui_language = 'en'  # Default to English

def t(key, *args):
    """
    Translation function - returns translated text for the current language.
    Args:
        key: Translation key
        *args: Format arguments for strings with placeholders
    """
    init_language_state()
    lang = st.session_state.ui_language
    
    # Get translation, fallback to English if not found
    translation = TRANSLATIONS.get(lang, {}).get(key, TRANSLATIONS["en"].get(key, key))
    
    # Format with arguments if provided
    if args:
        try:
            translation = translation.format(*args)
        except:
            pass  # If formatting fails, return unformatted string
    
    return translation

def render_language_selector():
    """Renders the language selector in the sidebar."""
    st.sidebar.markdown("---")
    
    languages = get_available_languages()
    current_lang = st.session_state.get('ui_language', 'en')
    
    selected_lang = st.sidebar.selectbox(
        t("language_selector"),
        options=list(languages.keys()),
        format_func=lambda x: f"{languages[x]}",
        index=list(languages.keys()).index(current_lang),
        key="language_selector"
    )
    
    if selected_lang != current_lang:
        st.session_state.ui_language = selected_lang
        st.rerun()
