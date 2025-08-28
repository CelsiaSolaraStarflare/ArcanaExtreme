import streamlit as st
import openai
from response import openai_api_call
import nltk

# Ensure NLTK data is available before importing NLTK functions
import arcana.nltk_setup
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
from arcana.fiber import FiberDBMS
from scripts.config import INDEX_FILE
import os
import json
import datetime
from docx import Document
from pptx import Presentation
import chardet
from PyPDF2 import PdfReader
import pandas as pd
from arcana.indexing import extract_keywords, detect_language

# NLTK data is now handled centrally in Arcanalte.py

# Chat History Management Functions

def get_chat_histories_dir():
    """Get the directory where chat histories are stored."""
    chat_dir = os.path.join(os.path.dirname(__file__), "chat_histories")
    os.makedirs(chat_dir, exist_ok=True)
    return chat_dir

def save_chat_history(session_name=None):
    """Save the current chat session to a file."""
    if "messages" not in st.session_state or not st.session_state.messages:
        st.warning("No messages to save!")
        return None
    
    if session_name is None:
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        session_name = f"chat_{timestamp}"
    
    # Filter out system messages for cleaner history
    user_messages = [msg for msg in st.session_state.messages if msg["role"] != "system"]
    
    chat_data = {
        "session_name": session_name,
        "timestamp": datetime.datetime.now().isoformat(),
        "messages": user_messages,
        "processed_file_name": st.session_state.get('processed_file_name', None)
    }
    
    chat_dir = get_chat_histories_dir()
    file_path = os.path.join(chat_dir, f"{session_name}.json")
    
    try:
        with open(file_path, 'w', encoding='utf-8') as f:
            json.dump(chat_data, f, indent=2, ensure_ascii=False)
        return file_path
    except Exception as e:
        st.error(f"Failed to save chat history: {e}")
        return None

def load_chat_history(file_path):
    """Load a chat session from a file."""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            chat_data = json.load(f)
        
        # Clear current messages completely and load from history
        st.session_state.messages = []
        
        # Add the initial system message first
        st.session_state.messages.append({
            "role": "system", 
            "content": "You are a helpful AI assistant named Arcana. You will be provided with search results from a user's documents. Your task is to answer the user's questions based *only* on the provided text. Cite the source document's name for all information you provide, like this: `(Source: document_name.pdf)`. If the provided text does not contain the answer, state that clearly. Be friendly, cute, and helpful."
        })
        
        # Add the welcome message only if there are no user messages in history
        user_messages_in_history = [msg for msg in chat_data.get("messages", []) if msg["role"] == "user"]
        if not user_messages_in_history:
            st.session_state.messages.append({
                "role": "assistant", 
                "content": "Hey, I'm Arcana, your Indexademics AI assistant. Ask me anything about your indexed files!"
            })
        
        # Load the chat history messages
        if chat_data.get("messages"):
            st.session_state.messages.extend(chat_data["messages"])
        
        # Restore processed file context if it exists
        st.session_state.processed_file_name = chat_data.get('processed_file_name', None)
        
        # Set the current session name to the loaded chat's name for continuous saving
        st.session_state.current_session_name = chat_data.get('session_name', None)
        
        return True
    except Exception as e:
        st.error(f"Failed to load chat history: {e}")
        return False

def get_available_chat_histories():
    """Get a list of available chat history files."""
    chat_dir = get_chat_histories_dir()
    histories = []
    
    for filename in os.listdir(chat_dir):
        if filename.endswith('.json'):
            file_path = os.path.join(chat_dir, filename)
            try:
                # Get file modification time for better sorting
                file_mtime = os.path.getmtime(file_path)
                
                with open(file_path, 'r', encoding='utf-8') as f:
                    chat_data = json.load(f)
                histories.append({
                    'filename': filename,
                    'filepath': file_path,
                    'session_name': chat_data.get('session_name', filename[:-5]),
                    'timestamp': chat_data.get('timestamp', 'Unknown'),
                    'message_count': len(chat_data.get('messages', [])),
                    'modified_time': file_mtime
                })
            except:
                continue
    
    # Sort by file modification time (most recently modified first)
    histories.sort(key=lambda x: x['modified_time'], reverse=True)
    return histories

def auto_generate_chat_title(messages):
    """Generate a meaningful title for the chat based on the conversation content."""
    try:
        # Extract user messages for title generation
        user_messages = [msg['content'] for msg in messages if msg['role'] == 'user']
        
        if not user_messages:
            return f"chat_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}"
        
        # Use first few user messages to generate title
        conversation_sample = " ".join(user_messages[:3])  # First 3 user messages
        
        # Limit the sample to avoid token limits
        if len(conversation_sample) > 500:
            conversation_sample = conversation_sample[:500] + "..."
        
        title_prompt = [
            {"role": "system", "content": "You are a helpful assistant that generates short, descriptive titles for chat conversations. Generate a concise title (3-6 words) that captures the main topic of the conversation. Only return the title, nothing else."},
            {"role": "user", "content": f"Generate a short title for this conversation: {conversation_sample}"}
        ]
        
        # Import here to avoid circular imports
        from response import openai_api_call
        
        title_generator = openai_api_call(title_prompt, "Normal")
        generated_title = "".join(title_generator).strip()
        
        # Clean up the title
        generated_title = generated_title.replace('"', '').replace("'", "")
        if len(generated_title) > 50:
            generated_title = generated_title[:50].strip()
        
        # Fallback if generation fails or is empty
        if not generated_title or len(generated_title) < 3:
            return f"chat_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}"
        
        return generated_title
        
    except Exception as e:
        print(f"Failed to generate title: {e}")
        return f"chat_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}"

def auto_save_current_chat():
    """Automatically save the current chat with an AI-generated title."""
    if "messages" not in st.session_state or not st.session_state.messages:
        return None
    
    # Check if there are any meaningful messages (user/assistant, not just system)
    meaningful_messages = [msg for msg in st.session_state.messages if msg["role"] in ["user", "assistant"]]
    
    if len(meaningful_messages) < 2:  # Need at least 1 user message and 1 assistant response
        return None
    
    with st.spinner("ü§ñ Generating title and saving chat..."):
        # Generate AI title
        auto_title = auto_generate_chat_title(st.session_state.messages)
        
        # Save with the generated title
        file_path = save_chat_history(auto_title)
        
        if file_path:
            st.success(f"üíæ Auto-saved chat as: '{auto_title}'")
            return file_path
        
    return None

def continuous_save_chat():
    """Continuously update the current chat session without showing notifications."""
    if "messages" not in st.session_state or not st.session_state.messages:
        return None
    
    # Check if there are any meaningful messages
    meaningful_messages = [msg for msg in st.session_state.messages if msg["role"] in ["user", "assistant"]]
    
    if len(meaningful_messages) < 1:  # Need at least some conversation
        return None
    
    # Check if we already have a current session name, if not generate one
    if not hasattr(st.session_state, 'current_session_name') or not st.session_state.current_session_name:
        if len(meaningful_messages) >= 2:  # Only generate title when we have full exchange
            st.session_state.current_session_name = auto_generate_chat_title(st.session_state.messages)
        else:
            st.session_state.current_session_name = f"chat_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}"
    
    # Save silently in the background
    file_path = save_chat_history(st.session_state.current_session_name)
    return file_path

def delete_chat_history(file_path):
    """Delete a chat history file."""
    try:
        os.remove(file_path)
        return True
    except Exception as e:
        st.error(f"Failed to delete chat history: {e}")
        return False

def extract_content_from_file(uploaded_file):
    """Extracts text content from an uploaded file object."""
    file_extension = os.path.splitext(uploaded_file.name)[1].lower()
    content = ""
    try:
        if file_extension == ".txt":
            # Use chardet to detect encoding for robust text file reading
            raw_data = uploaded_file.getvalue()
            encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'
            content = raw_data.decode(encoding, errors='replace')
        elif file_extension == ".docx":
            doc = Document(uploaded_file)
            content = "\n".join([para.text for para in doc.paragraphs])
        elif file_extension == ".pptx":
            presentation = Presentation(uploaded_file)
            all_texts = []
            for slide in presentation.slides:
                slide_texts = []
                if slide.shapes.title:
                    slide_texts.append(slide.shapes.title.text)
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_texts.append(shape.text_frame.text)  # type: ignore
                all_texts.append("\n".join(slide_texts))
            content = "\n".join(all_texts)
        elif file_extension == ".pdf":
            reader = PdfReader(uploaded_file)
            content = "\n".join([page.extract_text() or '' for page in reader.pages])
        elif file_extension in [".xls", ".xlsx", ".csv"]:
            # For simplicity, we'll read CSV from the buffer; Excel might need saving to disk
            if "xls" in file_extension:
                # To handle xls/xlsx from stream, pandas may need the 'openpyxl' or 'xlrd' engine
                df = pd.read_excel(uploaded_file, engine='openpyxl' if file_extension == ".xlsx" else 'xlrd')
            else:
                df = pd.read_csv(uploaded_file)
            content = df.to_string()
        else:
            st.warning(f"Unsupported file format: {file_extension}. Cannot interpret this file.")
            return None
    except Exception as e:
        st.error(f"Failed to process and interpret file {uploaded_file.name}: {e}")
        return None
    return content

def chatbot_page():
    st.title("Chat With Arcana")
    
    # Add custom CSS for ChatGPT-like styling
    st.markdown("""
    <style>
    .sidebar .block-container {
        padding-top: 1rem;
        padding-bottom: 1rem;
    }
    
    .chat-history-item {
        background-color: #f7f7f8;
        border-radius: 8px;
        padding: 8px 12px;
        margin: 4px 0;
        border: 1px solid #e5e5e7;
    }
    
    .chat-history-item:hover {
        background-color: #ececf1;
    }
    
    .new-chat-btn {
        background: linear-gradient(90deg, #1f2937 0%, #374151 100%);
        color: white;
        border: none;
        border-radius: 8px;
        padding: 12px;
        font-weight: 600;
        margin-bottom: 16px;
    }
    
    /* Reduce sidebar spacing */
    .css-1d391kg {
        padding-top: 1rem;
    }
    
    /* Style section headers */
    .sidebar-section-header {
        font-size: 14px;
        font-weight: 600;
        color: #374151;
        margin-bottom: 8px;
    }
    </style>
    """, unsafe_allow_html=True)

    # Initialize or load the database automatically
    if 'dbms' not in st.session_state or not isinstance(st.session_state.dbms, FiberDBMS):
        dbms = FiberDBMS()
        if os.path.exists(INDEX_FILE):
            with st.spinner("Loading existing database..."):
                try:
                    dbms.load_from_file(INDEX_FILE)
                    st.success("Database loaded successfully!")
                except Exception as e:
                    st.warning(f"Failed to load existing database: {e}. Starting with empty database.")
        else:
            st.info("No indexed files found. You can upload files directly or go to the 'Files' page to index documents.")
        st.session_state.dbms = dbms

    dbms = st.session_state.dbms
    
    with st.sidebar:
        # New Chat Button (prominent like ChatGPT)
        if st.button("‚ûï New Chat", key="new_chat_btn", use_container_width=True, type="primary"):
            # Auto-save current chat before clearing
            auto_save_current_chat()
            
            # Reset messages but keep the initial system prompt
            init_messages()
            st.session_state.processed_file_name = None
            # Reset the current session name for a fresh start
            st.session_state.current_session_name = None
            st.rerun()
        
        st.markdown("---")
        
        # Chat History Section (collapsible like ChatGPT)
        histories = get_available_chat_histories()
        if histories:
            with st.expander("üí¨ Chat History", expanded=True):
                # Show recent chats with clean design
                for i, history in enumerate(histories[:5]):  # Show only 5 most recent
                    with st.container():
                        col1, col2 = st.columns([4, 1])
                        with col1:
                            if st.button(
                                f"üìÑ {history['session_name'][:25]}...",
                                key=f"chat_{i}",
                                help=f"{history['message_count']} messages ‚Ä¢ {history['timestamp'][:10]}",
                                use_container_width=True
                            ):
                                if load_chat_history(history['filepath']):
                                    st.success("Chat loaded!")
                                    st.rerun()
                        
                        with col2:
                            if st.button("üóëÔ∏è", key=f"del_{i}", help="Delete chat"):
                                if delete_chat_history(history['filepath']):
                                    st.success("Deleted!")
                                    st.rerun()
                
                # Show more button if there are more than 5 chats
                if len(histories) > 5:
                    with st.expander("üìö Show All Chats", expanded=False):
                        for i, history in enumerate(histories[5:], start=5):
                            col1, col2 = st.columns([4, 1])
                            with col1:
                                if st.button(
                                    f"üìÑ {history['session_name'][:20]}...",
                                    key=f"chat_all_{i}",
                                    help=f"{history['message_count']} messages ‚Ä¢ {history['timestamp'][:10]}",
                                    use_container_width=True
                                ):
                                    if load_chat_history(history['filepath']):
                                        st.success("Chat loaded!")
                                        st.rerun()
                            with col2:
                                if st.button("üóëÔ∏è", key=f"del_all_{i}", help="Delete chat"):
                                    if delete_chat_history(history['filepath']):
                                        st.success("Deleted!")
                                        st.rerun()
        else:
            st.caption("üí¨ No chat history yet")
        
        st.markdown("---")
        
        # File Upload Section (prominent and visible)
        st.markdown("### üìÑ Upload Document")
        st.markdown("Upload a file to chat with it directly:")
        
        uploaded_file = st.file_uploader(
            "Choose a file to analyze",
            type=['txt', 'pdf', 'docx', 'pptx', 'csv', 'xlsx', 'xls'],
            help="Supported formats: PDF, Word, PowerPoint, Text, Excel, CSV"
        )
        
        # Show current file context if one is loaded
        if st.session_state.get('processed_file_name'):
            st.success(f"üìÑ Currently chatting with: **{st.session_state.processed_file_name}**")
            if st.button("‚ùå Clear File Context", use_container_width=True):
                st.session_state.processed_file_name = None
                # Remove file context from messages
                st.session_state.messages = [msg for msg in st.session_state.messages 
                                           if not (msg.get("role") == "system" and "uploaded the file" in msg.get("content", ""))]
                st.rerun()
        
        st.markdown("---")
        
        # Session Info (compact)
        if "messages" in st.session_state and st.session_state.messages:
            meaningful_messages = [msg for msg in st.session_state.messages if msg["role"] in ["user", "assistant"]]
            if len(meaningful_messages) >= 2:
                st.caption(f"üí¨ {len(meaningful_messages)} messages ‚Ä¢ Auto-saves when starting new chat")
        if uploaded_file is not None:
            # Check if this file has been processed already to avoid reprocessing on every rerun
            if st.session_state.get('processed_file_name') != uploaded_file.name:
                with st.spinner(f"üîç Processing {uploaded_file.name}..."):
                    file_content = extract_content_from_file(uploaded_file)
                    if file_content:
                        # Save the uploaded file to IDXDB/Uploads directory
                        uploads_dir = os.path.join(os.path.dirname(__file__), "IDXDB", "Uploads")
                        os.makedirs(uploads_dir, exist_ok=True)
                        
                        try:
                            # Save the original file
                            file_path = os.path.join(uploads_dir, uploaded_file.name)
                            with open(file_path, "wb") as f:
                                f.write(uploaded_file.getbuffer())
                            
                            # Also save extracted content as txt for easy reference
                            txt_filename = os.path.splitext(uploaded_file.name)[0] + "_extracted.txt"
                            txt_path = os.path.join(uploads_dir, txt_filename)
                            with open(txt_path, "w", encoding="utf-8") as f:
                                f.write(file_content)
                            
                            st.info(f"üíæ Saved to IDXDB/Uploads: {uploaded_file.name} + extracted text")
                        except Exception as e:
                            st.warning(f"Could not save file to IDXDB/Uploads: {e}")
                            # Continue with processing even if file saving fails
                        
                        # Index the new file content into the database
                        with st.spinner(f"üìö Indexing content..."):
                            lines = file_content.split('\n')
                            for line in lines:
                                line = line.strip()
                                if line:
                                    lang = detect_language(line)
                                    keywords = extract_keywords(line, lang)
                                    dbms.add_entry(name=uploaded_file.name, content=line, tags=keywords)
                            dbms.save(INDEX_FILE) # Save the updated index

                        # Add the file content as a system message for context, with priority instructions
                        context_message = (
                            f"The user has uploaded the file `{uploaded_file.name}`. "
                            f"For the user's next questions, you MUST prioritize the content of this file as the primary and ONLY source of information. "
                            f"Ignore any previous search results from the general document database. "
                            f"All answers must come directly from the file content provided below. "
                            f"Explicitly mention that you are answering based on the uploaded file."
                            f"\n\n--- FILE CONTENT ---\n{file_content}\n--- END FILE CONTENT ---"
                        )
                        st.session_state.messages.append({"role": "system", "content": context_message})
                        st.session_state.processed_file_name = uploaded_file.name
                        st.success(f"‚úÖ {uploaded_file.name} indexed and ready for questions!")
                        st.rerun()
                    else:
                        st.session_state.processed_file_name = None # Reset if processing fails

    # Remove the old clear button since it's now in sidebar as "New Chat"

    if "messages" not in st.session_state or not st.session_state.messages:
        init_messages()

    # Display existing conversation (excluding system messages)
    for message in st.session_state.messages:
        if message["role"] != "system":
            with st.chat_message(message["role"]):
                st.markdown(message["content"])

    # User input area
    user_input = st.chat_input("Ask me anything about your documents...")

    # Response type selector (more compact)
    with st.container():
        col1, col2 = st.columns([3, 1])
        with col2:
            response_type = st.selectbox(
                "Mode",
                ["Normal", "IDX", "Math"],
                help="""
                **Normal**: General conversation with search context
                **IDX**: Strictly based on indexed files
                **Math**: Specialized for mathematical queries
                """,
                label_visibility="collapsed"
            )

    if user_input:
        st.session_state.messages.append({"role": "user", "content": user_input})
        with st.chat_message("user"):
            st.markdown(user_input)
        
        # If a file has NOT been processed, search the database for context.
        # If a file HAS been processed, its context is already in the messages, so we skip this.
        if st.session_state.get('processed_file_name') is None:
            with st.spinner("Searching for relevant information..."):
                stop_words = set(stopwords.words('english'))
                words = word_tokenize(user_input)
                keywords = [word for word in words if word.lower() not in stop_words and word.isalpha()]
                
                # Use the dbms instance from session state
                results = dbms.query(" ".join(keywords), top_n=min(20, max(1, len(keywords))))
                results = results[:5]

                assistant_reply = ""
                if results:
                    assistant_reply += "Here are the top results from your documents:\n\n"
                    for idx, result in enumerate(results, 1):
                        assistant_reply += f"**Result {idx} from `{result['name']}`:**\n"
                        assistant_reply += f"_{result['content']}_\n\n"
                else:
                    assistant_reply = "I couldn't find any specific information related to your query in the indexed documents."

                st.session_state.messages.append({"role": "system", "content": assistant_reply})

        with st.spinner("Arcana is thinking..."):
            try:
                # Make sure to initialize 'processed_file_name' if it doesn't exist
                if 'processed_file_name' not in st.session_state:
                    st.session_state.processed_file_name = None

                with st.chat_message("assistant"):
                    # Use st.write_stream to render the response in real-time
                    response_generator = openai_api_call(st.session_state.messages, response_type)
                    collected_chunks = []

                    def stream_and_collect():
                        for chunk in response_generator:
                            collected_chunks.append(chunk)
                            yield chunk

                    st.write_stream(stream_and_collect())
                    full_response = "".join(collected_chunks)
                
                # Append the full response to the message history
                st.session_state.messages.append({"role": "assistant", "content": full_response})
                
                # Continuously save the chat in the background
                continuous_save_chat()
                
                st.rerun()

            except Exception as e:
                st.error(f"An error occurred while communicating with the AI: {e}")

def init_messages():
    """Initializes or resets the chat message history in the session state."""
    st.session_state.messages = [
        {"role": "assistant", "content": "Hey, I'm Arcana, your Indexademics AI assistant. Ask me anything about your indexed files!"},
        {"role": "system", "content": "You are a helpful AI assistant named Arcana. You will be provided with search results from a user's documents. Your task is to answer the user's questions based *only* on the provided text. Cite the source document's name for all information you provide, like this: `(Source: document_name.pdf)`. If the provided text does not contain the answer, state that clearly. Be friendly, cute, and helpful."},
    ]
